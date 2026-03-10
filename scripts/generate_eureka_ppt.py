#!/usr/bin/env python3
"""
OMNIX Eureka Dubai GCC 2026 -- Professional PowerPoint Presentation
Visual design with charts, diagrams, and graphics. March 2026.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_DIR = os.path.join(BASE_DIR, "docs", "business", "pdfs")
LOGO_PATH  = os.path.join(BASE_DIR, "omnix_web", "public", "logo.png")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ── Color palette ─────────────────────────────────────────────────────────────
NAVY      = RGBColor(10,  18,  40)
BLUE      = RGBColor(0,   80,  160)
BLUE2     = RGBColor(0,   120, 210)
BLUE_LIGHT= RGBColor(220, 234, 252)
WHITE     = RGBColor(255, 255, 255)
DARK      = RGBColor(20,  20,  40)
GRAY      = RGBColor(140, 148, 165)
GRAY_LITE = RGBColor(245, 247, 252)
GREEN     = RGBColor(20,  140, 70)
RED_SOFT  = RGBColor(190, 40,  40)
GOLD      = RGBColor(200, 155, 0)
ORANGE    = RGBColor(210, 100, 0)

W = Inches(13.33)
H = Inches(7.5)


def new_slide(prs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    _bg(s, WHITE)
    return s


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_width=0):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        sh.line.color.rgb = line_rgb
        sh.line.width = Pt(line_width)
    else:
        sh.line.fill.background()
    return sh


def _txt(slide, text, x, y, w, h, size=16, bold=False, color=None,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or DARK
    return tb


def _bullets(slide, items, x, y, w, h, size=13, color=None, indent=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        prefix = "\u2022  " if indent and not item.startswith("  ") else ""
        run.text = prefix + item.lstrip()
        run.font.size = Pt(size)
        run.font.color.rgb = color or DARK


def _footer(slide, slide_num, total=10):
    _rect(slide, 0, H - Inches(0.32), W, Inches(0.32), NAVY)
    _txt(slide, "OMNIX \u2014 Decision Governance Infrastructure  |  Eureka Dubai GCC 2026  |  Confidential",
         Inches(0.3), H - Inches(0.3), Inches(10), Inches(0.28),
         size=8, color=RGBColor(160, 180, 210), align=PP_ALIGN.LEFT)
    _txt(slide, f"omnixquantum.net  |  {slide_num}/{total}",
         Inches(11.5), H - Inches(0.3), Inches(1.6), Inches(0.28),
         size=8, color=RGBColor(160, 180, 210), align=PP_ALIGN.RIGHT)


def _header(slide, number, title, subtitle=None):
    _rect(slide, 0, 0, W, Inches(1.05), BLUE)
    _rect(slide, 0, 0, Inches(0.85), Inches(1.05), NAVY)
    num_str = f"0{number}" if number < 10 else str(number)
    _txt(slide, num_str, Inches(0.04), Inches(0.15), Inches(0.78), Inches(0.72),
         size=32, bold=True, color=RGBColor(200, 220, 255), align=PP_ALIGN.CENTER)
    _txt(slide, title.upper(), Inches(1.0), Inches(0.12), Inches(10.5), Inches(0.55),
         size=22, bold=True, color=WHITE)
    if subtitle:
        _txt(slide, subtitle, Inches(1.0), Inches(0.65), Inches(11.5), Inches(0.35),
             size=12, color=RGBColor(175, 210, 255), italic=True)
    _footer(slide, number)


def _metric_card(slide, x, y, w, h, value, label, bg=BLUE, val_color=WHITE, lbl_color=None):
    _rect(slide, x, y, w, h, bg)
    _txt(slide, value, x, y + Inches(0.1), w, h * 0.55,
         size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    _txt(slide, label, x, y + h * 0.55, w, h * 0.4,
         size=10, color=lbl_color or RGBColor(190, 215, 250), align=PP_ALIGN.CENTER)


def _bar_chart(slide, x, y, w, h, categories, series_data, title=""):
    cd = ChartData()
    cd.categories = categories
    for name, values in series_data:
        cd.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd
    ).chart
    chart.has_legend = len(series_data) > 1
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
    plot = chart.plots[0]
    plot.gap_width = 60
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        colors = [BLUE, BLUE2, GREEN, GOLD]
        series.format.fill.fore_color.rgb = colors[i % len(colors)]
    va = chart.value_axis
    va.has_major_gridlines = True
    va.tick_labels.font.size = Pt(8)
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(8)
    return chart


def _donut_chart(slide, x, y, w, h, categories, values, title=""):
    cd = ChartData()
    cd.categories = categories
    cd.add_series("", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, w, h, cd
    ).chart
    chart.has_legend = True
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
    slice_colors = [BLUE, BLUE2, GREEN, GOLD, ORANGE, RGBColor(160, 40, 120)]
    for i, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = slice_colors[i % len(slice_colors)]
    chart.legend.font.size = Pt(9)
    return chart


# =============================================================================
def generate():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # ── SLIDE 1: COVER ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)

    # Gold accent bar left
    _rect(s, 0, 0, Inches(0.18), H, GOLD)

    # Logo top-left (black bg logo works on dark slide)
    try:
        s.shapes.add_picture(LOGO_PATH, Inches(0.35), Inches(0.3), Inches(2.8), Inches(2.8))
    except Exception:
        _txt(s, "OMNIX", Inches(0.35), Inches(0.5), Inches(2.8), Inches(1.2),
             size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Title block
    _txt(s, "OMNIX QUANTUM", Inches(3.5), Inches(0.6), Inches(9.4), Inches(1.2),
         size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _txt(s, "Decision Governance Infrastructure", Inches(3.5), Inches(1.75), Inches(9.4), Inches(0.65),
         size=22, color=RGBColor(180, 215, 255), align=PP_ALIGN.LEFT)
    _txt(s, "for Automated Financial Systems", Inches(3.5), Inches(2.35), Inches(9.4), Inches(0.5),
         size=16, color=GRAY, align=PP_ALIGN.LEFT, italic=True)

    # Divider line
    _rect(s, Inches(3.5), Inches(3.0), Inches(9.3), Inches(0.04), GOLD)

    # Key info
    info_items = [
        ("\u2605  SEMIFINALIST \u2014 Eureka Dubai GCC 2026", 17, True, GOLD),
        ("Pre-Seed Round  |  $500,000 USD  |  16.7% equity  |  $3M pre-money valuation", 13, False, RGBColor(190, 210, 245)),
        ("Founder & CEO: Harold Nunes   |   contacto@omnixquantum.net   |   omnixquantum.net", 11, False, GRAY),
        ("March 2026", 11, False, GRAY),
    ]
    y_pos = 3.15
    for txt, sz, bld, clr in info_items:
        _txt(s, txt, Inches(3.5), Inches(y_pos), Inches(9.3), Inches(0.5),
             size=sz, bold=bld, color=clr, align=PP_ALIGN.LEFT)
        y_pos += 0.48

    # Bottom tagline
    _rect(s, Inches(0.18), H - Inches(0.85), W - Inches(0.18), Inches(0.85), BLUE)
    _txt(s, '"Intelligence may propose. OMNIX decides whether it is allowed to act."',
         Inches(0.4), H - Inches(0.78), W - Inches(0.5), Inches(0.65),
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    # ── SLIDE 2: PROBLEM ─────────────────────────────────────────────────────
    s = new_slide(prs)
    _header(s, 2, "The Problem", "Automated systems fail \u2014 and they fail at machine speed")

    # Quote banner
    _rect(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.75), NAVY)
    _txt(s, "\u201cWhen automated systems fail, they do not fail slowly. They fail at machine speed.\u201d",
         Inches(0.6), Inches(1.22), Inches(12.1), Inches(0.6),
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    # 3 failure mode cards
    failures = [
        ("SIGNAL\nFAILURE", "Systems act on false signals during volatile or anomalous market conditions. No checkpoint stops a miscalibrated entry.", RGBColor(150, 20, 20)),
        ("COHERENCE\nFAILURE", "Multiple internal signals disagree, yet the system executes anyway. Internal conflict = guaranteed bad decision.", RGBColor(160, 80, 0)),
        ("TAIL RISK\nFAILURE", "Low-probability, high-impact events (Black Swans) are not detected until losses are already cascading.", RGBColor(80, 0, 130)),
    ]
    card_w = Inches(4.0)
    card_h = Inches(2.4)
    for i, (title, desc, color) in enumerate(failures):
        cx = Inches(0.4) + i * (card_w + Inches(0.15))
        _rect(s, cx, Inches(2.05), card_w, Inches(0.6), color)
        _txt(s, title, cx, Inches(2.08), card_w, Inches(0.55),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _rect(s, cx, Inches(2.65), card_w, card_h - Inches(0.6), GRAY_LITE)
        _txt(s, desc, cx + Inches(0.15), Inches(2.72), card_w - Inches(0.3), card_h - Inches(0.75),
             size=12, color=DARK, wrap=True)

    # Bottom facts
    _rect(s, Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.04), BLUE)
    facts = [
        "MiCA (EU 2025+): regulators now require decision accountability documentation",
        "ADGM \u0026 DIFC: AI governance documentation mandatory for market access in GCC",
        "Prop trading firms lose 15\u201340% of capital in first year of automated deployment",
        "Post-quantum security becoming a procurement requirement at institutional level",
    ]
    _bullets(s, facts, Inches(0.4), Inches(4.72), Inches(12.5), Inches(2.0), size=13)

    # ── SLIDE 3: SOLUTION ────────────────────────────────────────────────────
    s = new_slide(prs)
    _header(s, 3, "The Solution", "A fail-closed governance layer that authorizes every decision")

    # Visual architecture diagram
    box_h = Inches(0.85)
    box_y = Inches(1.2)
    # Layer 1 - Automated System
    _rect(s, Inches(0.5), box_y, Inches(3.5), box_h, GRAY_LITE,
          line_rgb=GRAY, line_width=1)
    _txt(s, "AUTOMATED SYSTEM", Inches(0.5), box_y + Inches(0.05), Inches(3.5), Inches(0.35),
         size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    _txt(s, "Signal generators, strategies,\nalgorithms, AI models",
         Inches(0.5), box_y + Inches(0.38), Inches(3.5), Inches(0.44),
         size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # Arrow right
    _rect(s, Inches(4.1), box_y + Inches(0.35), Inches(0.6), Inches(0.15), BLUE)
    _txt(s, "\u25ba", Inches(4.55), box_y + Inches(0.28), Inches(0.25), Inches(0.3),
         size=14, color=BLUE, align=PP_ALIGN.CENTER)

    # OMNIX box (center - prominent)
    _rect(s, Inches(4.75), box_y - Inches(0.2), Inches(4.0), box_h + Inches(0.4), BLUE)
    _txt(s, "OMNIX GOVERNANCE", Inches(4.75), box_y - Inches(0.12), Inches(4.0), Inches(0.45),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s, "8-CHECKPOINT VALIDATION ENGINE", Inches(4.75), box_y + Inches(0.3), Inches(4.0), Inches(0.3),
         size=10, color=RGBColor(190, 220, 255), align=PP_ALIGN.CENTER)
    _txt(s, "PQC-SIGNED AUDIT TRAIL", Inches(4.75), box_y + Inches(0.58), Inches(4.0), Inches(0.28),
         size=9, color=RGBColor(200, 220, 255), align=PP_ALIGN.CENTER, italic=True)

    # Arrow right (or block)
    _rect(s, Inches(8.85), box_y + Inches(0.35), Inches(0.6), Inches(0.15), BLUE)
    _txt(s, "\u25ba", Inches(9.3), box_y + Inches(0.28), Inches(0.25), Inches(0.3),
         size=14, color=BLUE, align=PP_ALIGN.CENTER)

    # Execution
    _rect(s, Inches(9.55), box_y, Inches(3.3), box_h, GRAY_LITE,
          line_rgb=GRAY, line_width=1)
    _txt(s, "EXECUTION", Inches(9.55), box_y + Inches(0.05), Inches(3.3), Inches(0.35),
         size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    _txt(s, "Only authorized decisions\nreach the market",
         Inches(9.55), box_y + Inches(0.38), Inches(3.3), Inches(0.44),
         size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # VETO arrow (down from OMNIX)
    _rect(s, Inches(6.4), box_y + box_h + Inches(0.2), Inches(0.12), Inches(0.5), RED_SOFT)
    _txt(s, "\u25bc  VETO (BLOCKED)", Inches(5.6), box_y + box_h + Inches(0.68), Inches(1.8), Inches(0.3),
         size=10, bold=True, color=RED_SOFT, align=PP_ALIGN.CENTER)

    _rect(s, Inches(0.4), Inches(2.85), Inches(12.5), Inches(0.04), BLUE_LIGHT)

    # Two columns
    col_w = Inches(6.1)
    _rect(s, Inches(0.4), Inches(2.95), col_w, Inches(0.32), GREEN)
    _txt(s, "OMNIX DOES \u2714", Inches(0.5), Inches(2.98), col_w, Inches(0.28),
         size=12, bold=True, color=WHITE)
    does = [
        "Intercepts every signal before execution",
        "Runs 8 independent validation checkpoints",
        "Blocks decisions that fail any single checkpoint",
        "Signs every decision with post-quantum cryptography",
        "Maintains a tamper-proof verifiable audit trail",
    ]
    _bullets(s, does, Inches(0.5), Inches(3.32), col_w - Inches(0.2), Inches(2.8),
             size=13, color=DARK)

    _rect(s, Inches(6.8), Inches(2.95), col_w, Inches(0.32), RED_SOFT)
    _txt(s, "OMNIX DOES NOT \u2718", Inches(6.9), Inches(2.98), col_w, Inches(0.28),
         size=12, bold=True, color=WHITE)
    does_not = [
        "Generate trading signals or strategies",
        "Compete with existing trading platforms",
        "Manage or hold client assets",
        "Require replacing existing infrastructure",
        "Optimize for returns \u2014 only for capital safety",
    ]
    _bullets(s, does_not, Inches(6.9), Inches(3.32), col_w - Inches(0.2), Inches(2.8),
             size=13, color=DARK)

    # ── SLIDE 4: HOW IT WORKS ────────────────────────────────────────────────
    s = new_slide(prs)
    _header(s, 4, "How It Works", "8 independent checkpoints \u2014 every decision must earn execution")

    checkpoints = [
        ("CP-0", "Signal Integrity\nValidator (SIV)", "Data quality gate\nRejects stale/corrupted data"),
        ("CP-1", "Monte Carlo\nVETO", "10,000 simulations\nBlocks negative EV trades"),
        ("CP-2", "RMS VETO", "Risk limits\nHard exposure caps"),
        ("CP-3", "VETO Early\nReturn", "Fast-fail gate\nCritical threshold breach"),
        ("CP-4", "Coherence\nEngine (6-tier)", "Signal agreement\n6-dimension analysis"),
        ("CP-5", "Adaptive\nCoherence Gate", "Dynamic threshold\nRegime-adjusted"),
        ("CP-7", "Temporal\nCoherence (TCV)", "Backward trajectory\nHistory validation"),
        ("CP-8", "Edge Confirmation\nWindow (ECW)", "2 consecutive cycles\nof confirmed edge"),
    ]

    box_w = Inches(1.5)
    box_h = Inches(1.5)
    gap   = Inches(0.12)
    start_x = Inches(0.35)
    cp_y  = Inches(1.2)

    cp_colors = [NAVY, BLUE, BLUE2, RGBColor(0, 140, 200),
                 RGBColor(0, 100, 175), RGBColor(0, 85, 155),
                 RGBColor(20, 65, 140), RGBColor(40, 55, 120)]

    for i, (cp_id, name, desc) in enumerate(checkpoints):
        cx = start_x + i * (box_w + gap)
        _rect(s, cx, cp_y, box_w, box_h, cp_colors[i])
        _txt(s, cp_id, cx, cp_y + Inches(0.06), box_w, Inches(0.3),
             size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        _txt(s, name, cx, cp_y + Inches(0.32), box_w, Inches(0.55),
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, desc, cx, cp_y + Inches(0.85), box_w, Inches(0.6),
             size=8, color=RGBColor(190, 215, 250), align=PP_ALIGN.CENTER)
        # Arrow (not after last)
        if i < len(checkpoints) - 1:
            ax = cx + box_w + Inches(0.02)
            _txt(s, "\u25ba", ax, cp_y + Inches(0.62), gap + Inches(0.05), Inches(0.3),
                 size=10, color=BLUE2, align=PP_ALIGN.CENTER)

    # EGL bar
    _rect(s, Inches(0.35), Inches(2.85), Inches(12.6), Inches(0.04), GOLD)
    _rect(s, Inches(0.35), Inches(2.93), Inches(12.6), Inches(0.65), NAVY)
    _txt(s, "EXIT GOVERNANCE LAYER (EGL)  \u2014  3 Gates:", Inches(0.55), Inches(2.97),
         Inches(4.5), Inches(0.55), size=11, bold=True, color=GOLD)
    _txt(s, "Gate 1: Regime-Adjusted Thresholds     \u25ba     Gate 2: Exit Coherence Gate     \u25ba     Gate 3: TCV Exit Check",
         Inches(5.0), Inches(3.0), Inches(7.9), Inches(0.5),
         size=11, color=RGBColor(190, 215, 250), align=PP_ALIGN.CENTER)

    # Fail-closed principle
    _rect(s, Inches(0.35), Inches(3.72), Inches(12.6), Inches(0.04), BLUE)
    _txt(s, "FAIL-CLOSED PRINCIPLE:", Inches(0.45), Inches(3.82),
         Inches(2.5), Inches(0.4), size=12, bold=True, color=BLUE)
    _txt(s, "Every checkpoint has independent veto authority. One failure = execution blocked. No exceptions.",
         Inches(2.95), Inches(3.82), Inches(10.0), Inches(0.4),
         size=12, color=DARK)

    # PQC note
    _txt(s, "[PQC]  Every decision \u2014 entry or exit \u2014 is cryptographically signed with NIST-standardized post-quantum algorithms (Dilithium-3) and stored in a tamper-proof public audit trail at omnixquantum.net/verify",
         Inches(0.45), Inches(4.32), Inches(12.5), Inches(0.55),
         size=11, color=GRAY, italic=True)

    # ── SLIDE 5: TRACTION ────────────────────────────────────────────────────
    s = new_slide(prs)
    _header(s, 5, "Traction", "Live system \u2014 verified governance metrics from production")

    # 4 metric cards
    metrics = [
        ("50,688", "Governance decisions"),
        ("728,868", "Shadow evaluations"),
        ("100%", "Veto accuracy (50 validated)"),
        ("Semifinalist", "Eureka Dubai GCC 2026"),
    ]
    card_w = Inches(2.95)
    for i, (val, lbl) in enumerate(metrics):
        cx = Inches(0.35) + i * (card_w + Inches(0.13))
        _rect(s, cx, Inches(1.18), card_w, Inches(0.95), BLUE)
        _txt(s, val, cx, Inches(1.22), card_w, Inches(0.55),
             size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, lbl, cx, Inches(1.72), card_w, Inches(0.38),
             size=10, color=RGBColor(190, 215, 250), align=PP_ALIGN.CENTER)

    # Bar chart: veto types
    _bar_chart(s,
               x=Inches(0.35), y=Inches(2.25), w=Inches(6.0), h=Inches(2.8),
               categories=["BLACK_SWAN", "ECW WAITING", "COHERENCE GATE"],
               series_data=[("Veto Count (Jan 15 - Mar 2026)", [396067, 162164, 121157])],
               title="Veto Distribution by Type (Jan 15, 2026 \u2013 present)")

    # Timeline
    _txt(s, "System timeline:", Inches(6.6), Inches(2.22), Inches(6.5), Inches(0.38),
         size=13, bold=True, color=BLUE)
    phases = [
        (RGBColor(140, 20, 20), "Phase 0 (Jul\u2013Aug 2025)", "1,115 real trades on Kraken\nReal capital, pre-governance baseline"),
        (RGBColor(160, 100, 0), "Learning Baseline (Nov 2025\u2013Jan 14, 2026)", "119 paper trades\nRisk engine calibrated, veto system tuned"),
        (GREEN, "Official Track Record (Jan 15, 2026\u2013present)", "679,388 veto decisions\n0 executions, capital fully protected"),
    ]
    py = Inches(2.65)
    for color, title, desc in phases:
        _rect(s, Inches(6.6), py, Inches(0.28), Inches(0.7), color)
        _txt(s, title, Inches(7.0), py + Inches(0.02), Inches(6.1), Inches(0.32),
             size=11, bold=True, color=DARK)
        _txt(s, desc, Inches(7.0), py + Inches(0.3), Inches(6.1), Inches(0.38),
             size=10, color=GRAY, italic=True)
        py += Inches(0.82)

    _rect(s, Inches(0.35), Inches(5.15), Inches(12.6), Inches(0.04), BLUE_LIGHT)
    _txt(s, "\u26a0  Phase 0 = real capital on Kraken exchange (Jul\u2013Aug 2025). Learning Baseline & Official Track Record = paper/simulated trading. Never mixed in reporting.",
         Inches(0.4), Inches(5.25), Inches(12.5), Inches(0.45),
         size=10, color=GRAY, italic=True)

    # ── SLIDE 6: MARKET OPPORTUNITY ──────────────────────────────────────────
    s = new_slide(prs)
    _header(s, 6, "Market Opportunity", "$5.4B TAM \u2014 compliance-driven demand accelerating in 2025+")

    # Donut chart TAM/SAM
    _donut_chart(s,
                 x=Inches(0.3), y=Inches(1.15), w=Inches(5.5), h=Inches(4.2),
                 categories=["SAM (Serviceable, Y1-3)", "Remaining TAM", "Emerging verticals (Y2+)"],
                 values=[540, 4860, 12400],
                 title="Market Size ($M USD)")

    # Right side: Market breakdown
    _txt(s, "Market drivers:", Inches(6.0), Inches(1.18), Inches(7.0), Inches(0.38),
         size=13, bold=True, color=BLUE)
    drivers = [
        "MiCA (EU 2025+) \u2014 mandates decision accountability for all automated crypto systems",
        "ADGM & DIFC (UAE) \u2014 AI governance documentation required for market access",
        "Sharia compliance \u2014 GCC institutions require validated decision documentation for audit",
        "Post-quantum mandate \u2014 regulated GCC institutions beginning to require PQC procurement",
        "Automation growing faster than governance frameworks \u2014 the gap widens each year",
    ]
    _bullets(s, drivers, Inches(6.0), Inches(1.62), Inches(7.1), Inches(2.2), size=12)

    # Segment cards
    segs = [
        ("~4,000", "Prop trading firms globally", BLUE),
        ("~500", "Regulated exchanges", BLUE2),
        ("~2,000", "Hedge funds (relevant)", GREEN),
        ("~18,000", "Total target clients", NAVY),
    ]
    sw = Inches(1.68)
    sx = Inches(6.0)
    for i, (n, label, col) in enumerate(segs):
        bx = sx + i * (sw + Inches(0.06))
        _rect(s, bx, Inches(4.0), sw, Inches(0.85), col)
        _txt(s, n, bx, Inches(4.05), sw, Inches(0.45),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, label, bx, Inches(4.48), sw, Inches(0.34),
             size=9, color=RGBColor(200, 220, 255), align=PP_ALIGN.CENTER)

    # Key number row
    _rect(s, Inches(0.35), Inches(5.22), Inches(12.6), Inches(0.55), NAVY)
    _txt(s, "TAM: $5.4B USD (18,000 \xd7 $300K avg/yr)    \u2502    SAM: $540M (top 10% with active compliance pressure)    \u2502    First year target: 3 enterprise pilots",
         Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
         size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # ── SLIDE 7: BUSINESS MODEL ───────────────────────────────────────────────
    s = new_slide(prs)
    _header(s, 7, "Business Model & Projections", "Recurring SaaS + API-first + institutional licensing")

    # 4 pricing tier cards
    tiers = [
        ("Enterprise\nPilot", "$10K\u2013$15K\n/month", "Prop firms, exchanges", NAVY),
        ("Enterprise\nFull", "$20K\u2013$35K\n/month", "Hedge funds, reg. funds", BLUE),
        ("API\nPlatform", "$0.01\u2013$0.05\n/call", "High-volume platforms", BLUE2),
        ("B2C\nAdvanced", "$149\u2013$499\n/month", "Advanced retail traders", RGBColor(0, 140, 200)),
    ]
    tw = Inches(3.05)
    for i, (name, price, target, col) in enumerate(tiers):
        tx = Inches(0.35) + i * (tw + Inches(0.1))
        _rect(s, tx, Inches(1.18), tw, Inches(1.05), col)
        _txt(s, name, tx, Inches(1.2), tw, Inches(0.38),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, price, tx, Inches(1.56), tw, Inches(0.5),
             size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        _txt(s, target, tx, Inches(2.03), tw, Inches(0.18),
             size=9, color=GRAY_LITE, align=PP_ALIGN.CENTER, italic=True)

    # Revenue bar chart
    _bar_chart(s,
               x=Inches(0.35), y=Inches(2.35), w=Inches(7.5), h=Inches(3.0),
               categories=["Y1 (2026)", "Y2 (2027)", "Y3 (2028)", "Y4 (2029)", "Y5 (2030)"],
               series_data=[("Revenue ($USD)", [300000, 1200000, 4800000, 13000000, 26000000])],
               title="Revenue Projections Y1\u2013Y5")

    # MOIC cards
    moic_data = [
        ("14.7x", "Conservative\nMOIC", RGBColor(0, 100, 60)),
        ("41x", "Base\nMOIC", BLUE),
        ("102x", "Optimistic\nMOIC", NAVY),
    ]
    mw = Inches(1.72)
    for i, (val, lbl, col) in enumerate(moic_data):
        mx = Inches(8.1) + i * (mw + Inches(0.13))
        _rect(s, mx, Inches(2.35), mw, Inches(0.88), col)
        _txt(s, val, mx, Inches(2.38), mw, Inches(0.52),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, lbl, mx, Inches(2.88), mw, Inches(0.32),
             size=9, color=RGBColor(200, 220, 255), align=PP_ALIGN.CENTER)

    _txt(s, "Investor return scenarios\n(pre-seed entry):", Inches(8.1), Inches(3.4), Inches(5.1), Inches(0.6),
         size=11, bold=True, color=BLUE)
    _bullets(s, [
        "Break-even target: Q4 2026",
        "First revenue: Q2-Q3 2026",
        "Series A ready: Year 3",
        "Path to $26M ARR by Y5",
    ], Inches(8.1), Inches(4.05), Inches(5.1), Inches(1.4), size=12)

    # ── SLIDE 8: COMPETITIVE ADVANTAGE ────────────────────────────────────────
    s = new_slide(prs)
    _header(s, 8, "Competitive Advantage", "The only fail-closed governance layer with PQC audit trail")

    headers = ["Feature", "Retail\nBots", "Risk\nPlatforms", "Quant\nFunds", "OMNIX"]
    col_ws  = [Inches(4.5), Inches(1.85), Inches(1.85), Inches(1.85), Inches(2.1)]
    rows = [
        ["Fail-closed architecture",          "No", "Partial", "No",      "YES \u2713"],
        ["8 independent checkpoints",          "No", "No",      "No",      "YES \u2713"],
        ["Post-quantum cryptographic trail",   "No", "No",      "No",      "YES \u2713"],
        ["Public verifiable receipt URL",      "No", "No",      "No",      "YES \u2713"],
        ["MiCA / ADGM compliance-ready",       "No", "Partial", "No",      "YES \u2713"],
        ["Shadow portfolio engine",            "No", "No",      "Internal","YES \u2713"],
        ["Real-time Black Swan detection",     "No", "Partial", "Partial", "YES \u2713"],
    ]

    cell_colors = {
        "YES \u2713": GREEN,
        "No": RGBColor(200, 50, 50),
        "Partial": ORANGE,
        "Internal": RGBColor(140, 100, 0),
    }

    row_h  = Inches(0.52)
    table_y = Inches(1.18)
    start_x = Inches(0.35)

    # Header row
    _rect(s, start_x, table_y, sum(col_ws), row_h, NAVY)
    x = start_x
    for ci, h in enumerate(headers):
        _txt(s, h, x + Inches(0.06), table_y + Inches(0.06), col_ws[ci] - Inches(0.1), row_h,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        x += col_ws[ci]

    for ri, row in enumerate(rows):
        ry = table_y + (ri + 1) * row_h
        bg_col = GRAY_LITE if ri % 2 == 0 else WHITE
        _rect(s, start_x, ry, sum(col_ws), row_h, bg_col)
        x = start_x
        for ci, cell in enumerate(row):
            if ci == 0:
                _txt(s, cell, x + Inches(0.1), ry + Inches(0.1),
                     col_ws[ci] - Inches(0.15), row_h - Inches(0.1),
                     size=12, color=DARK)
            else:
                c_bg = cell_colors.get(cell, GRAY_LITE)
                _rect(s, x + Inches(0.04), ry + Inches(0.06),
                      col_ws[ci] - Inches(0.08), row_h - Inches(0.12), c_bg)
                _txt(s, cell, x + Inches(0.04), ry + Inches(0.1),
                     col_ws[ci] - Inches(0.08), row_h - Inches(0.15),
                     size=11, bold=(cell == "YES \u2713"), color=WHITE, align=PP_ALIGN.CENTER)
            x += col_ws[ci]

    _rect(s, Inches(0.35), Inches(5.12), Inches(12.6), Inches(0.55), BLUE)
    _txt(s, "The fundamental difference: OMNIX is fail-CLOSED. Every other system is fail-OPEN. A decision must earn execution.",
         Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── SLIDE 9: EXPERT VALIDATION ────────────────────────────────────────────
    s = new_slide(prs)
    _header(s, 9, "Expert Validation", "6 domain leaders independently validated the architecture via LinkedIn")

    experts = [
        ("JM", "James Moore", "Founder/CEO, Nova Jema AI Systems\nAI governance for healthcare & public infra",
         "\"Harold, this is exactly the layer I've been trying to articulate... the boundary between recommendation and binding authority.\"",
         BLUE),
        ("WF", "William Fedorich", "Author & AI Risk Analyst\nPublisher, The Bill Fedorich Letter",
         "\"Ensuring every AI decision is vetted through multiple checkpoints is critical for high-stakes environments like trading and insurance.\"",
         RGBColor(0, 100, 80)),
        ("CT", "Christopher Turk", "Quantum Security Architecture\nZero Trust & Agentic AI Security",
         "\"This is honestly pretty good.\" Extended technical exchange validating Dilithium-3 selection and PQC provenance model.",
         NAVY),
        ("GY", "Guomin Yang", "Strategist & Thought Leader\nCreator, RDL Framework \u2014 AI Governance",
         "Independently cited Harold by name in a published article, mapping OMNIX to his 3-layer governance model (Policy \u2192 Operational \u2192 Structural).",
         RGBColor(100, 0, 130)),
        ("HW", "HIL-AIW", "AI Governance Platform\n925+ LinkedIn followers",
         "\"Your 8-checkpoint system with Monte Carlo VETO is brilliant \u2014 most teams aren't thinking about quantum-resistant auditability yet.\"",
         RGBColor(0, 120, 160)),
        ("MM", "Mostafa Monsour", "Founder, ULTRA MATRIX\nCognitive & Strategic Architect",
         "\"Your work with OMNIX Quantum appears to be exploring an important part of that stack.\"",
         RGBColor(150, 80, 0)),
    ]

    card_w = Inches(4.1)
    card_h = Inches(1.72)
    gap_x  = Inches(0.2)
    gap_y  = Inches(0.12)

    for i, (initials, name, role, quote, color) in enumerate(experts):
        col = i % 3
        row = i // 3
        cx = Inches(0.35) + col * (card_w + gap_x)
        cy = Inches(1.18) + row * (card_h + gap_y)

        _rect(s, cx, cy, card_w, card_h, GRAY_LITE)
        # Avatar circle (drawn as colored rect)
        _rect(s, cx + Inches(0.12), cy + Inches(0.12), Inches(0.6), Inches(0.6), color)
        _txt(s, initials, cx + Inches(0.12), cy + Inches(0.15), Inches(0.6), Inches(0.5),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Name + role
        _txt(s, name, cx + Inches(0.82), cy + Inches(0.1), card_w - Inches(0.95), Inches(0.32),
             size=11, bold=True, color=DARK)
        _txt(s, role, cx + Inches(0.82), cy + Inches(0.4), card_w - Inches(0.95), Inches(0.38),
             size=9, color=GRAY, italic=True)
        # Quote
        _rect(s, cx + Inches(0.1), cy + Inches(0.85), card_w - Inches(0.2), Inches(0.04), color)
        _txt(s, quote, cx + Inches(0.12), cy + Inches(0.92), card_w - Inches(0.25), Inches(0.75),
             size=9, color=DARK, italic=True)

    # ── SLIDE 10: THE ASK + CLOSE ─────────────────────────────────────────────
    s = new_slide(prs)
    _header(s, 10, "The Investment", "$500K pre-seed \u2014 governance infrastructure for a $5.4B market")

    # 4 metric cards
    asks = [
        ("$500K", "Raise target"),
        ("16.7%", "Equity offered"),
        ("$3M", "Pre-money valuation"),
        ("Q4 2026", "Break-even target"),
    ]
    aw = Inches(2.95)
    for i, (v, l) in enumerate(asks):
        ax = Inches(0.35) + i * (aw + Inches(0.13))
        _rect(s, ax, Inches(1.18), aw, Inches(0.95), BLUE)
        _txt(s, v, ax, Inches(1.22), aw, Inches(0.55),
             size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, l, ax, Inches(1.72), aw, Inches(0.38),
             size=10, color=RGBColor(190, 215, 250), align=PP_ALIGN.CENTER)

    # Donut chart: use of funds
    _donut_chart(s,
                 x=Inches(0.35), y=Inches(2.25), w=Inches(5.0), h=Inches(3.0),
                 categories=["Infrastructure & Security (40%)", "Sales & Business Dev (35%)", "Team & Operations (25%)"],
                 values=[40, 35, 25],
                 title="Use of Funds")

    # Right: fund use details + MOIC
    fund_details = [
        (BLUE, "40% Infrastructure & Security ($200K)", "Production hardening, PQC security audit, compliance certifications"),
        (GREEN, "35% Sales & Business Development ($175K)", "Institutional pilots in US & GCC, ADGM setup, Sharia compliance assessment"),
        (RGBColor(0, 140, 200), "25% Team & Operations ($125K)", "First technical hire, legal entity, operating runway through Q4 2026"),
    ]
    fy = Inches(2.3)
    for color, title, desc in fund_details:
        _rect(s, Inches(5.6), fy, Inches(0.22), Inches(0.68), color)
        _txt(s, title, Inches(5.95), fy + Inches(0.02), Inches(7.2), Inches(0.3),
             size=11, bold=True, color=DARK)
        _txt(s, desc, Inches(5.95), fy + Inches(0.3), Inches(7.2), Inches(0.35),
             size=10, color=GRAY)
        fy += Inches(0.82)

    # MOIC bar
    _rect(s, Inches(5.6), Inches(4.7), Inches(7.5), Inches(0.35), NAVY)
    _txt(s, "MOIC at exit:", Inches(5.7), Inches(4.74), Inches(1.5), Inches(0.28),
         size=11, bold=True, color=GOLD)
    moic_items = [("14.7x", "Conservative", GREEN), ("41x", "Base", BLUE), ("102x", "Optimistic", NAVY)]
    for i, (val, label, col) in enumerate(moic_items):
        mx = Inches(7.4) + i * Inches(1.85)
        _txt(s, f"{val} ({label})", mx, Inches(4.74), Inches(1.8), Inches(0.28),
             size=11, bold=True, color=GOLD if i == 2 else WHITE, align=PP_ALIGN.CENTER)

    # Close banner
    _bg(s, WHITE)
    _rect(s, 0, Inches(5.22), W, Inches(1.62), NAVY)
    _txt(s, "The governance layer is missing. OMNIX is building it.",
         Inches(0), Inches(5.35), W, Inches(0.6),
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s, "Harold Nunes \u2014 Founder & CEO  |  contacto@omnixquantum.net  |  omnixquantum.net  |  omnixquantum.net/verify",
         Inches(0), Inches(5.95), W, Inches(0.4),
         size=12, color=RGBColor(180, 210, 255), align=PP_ALIGN.CENTER)
    _txt(s, "\u2605  SEMIFINALIST \u2014 Eureka Dubai GCC 2026  \u2605",
         Inches(0), Inches(6.38), W, Inches(0.38),
         size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    out = os.path.join(OUTPUT_DIR, "OMNIX_Eureka_Presentation.pptx")
    prs.save(out)
    size_kb = os.path.getsize(out) / 1024
    print(f"Generated: {out}")
    print(f"File size: {size_kb:.1f} KB")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    generate()
